"""Generate a PowerPoint from `context.md` using slide_bg images as slide backgrounds.
"""

from pathlib import Path
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

try:
    from PIL import Image

    PIL_AVAILABLE = True
except Exception:
    PIL_AVAILABLE = False


ROOT = Path(__file__).resolve().parent
CONTEXT = ROOT / "context.md"
OUT_FILE = ROOT / "output.pptx"


def parse_hex_color(s: str):
    """Parse a hex color like '#fff' or '#ffffff' into an RGBColor or return None."""
    if not s:
        return None
    t = s.strip()
    if t.startswith("#"):
        t = t[1:]
    if len(t) == 3:
        t = "".join([c * 2 for c in t])
    if len(t) != 6:
        return None
    try:
        r = int(t[0:2], 16)
        g = int(t[2:4], 16)
        b = int(t[4:6], 16)
        return RGBColor(r, g, b)
    except Exception:
        return None


def parse_blocks(text: str):
    parts = [p.strip() for p in re.split(r"\n---\n", text) if p.strip()]
    blocks = []
    for part in parts:
        lines = part.splitlines()
        data = {}
        i = 0
        while i < len(lines):
            line = lines[i].rstrip()
            if not line or line.startswith("##"):
                i += 1
                continue
            if ":" in line:
                key, val = line.split(":", 1)
                key = key.strip()
                val = val.strip()
                if val == "":
                    # gather indented block
                    i += 1
                    collected = []
                    while i < len(lines) and (
                        lines[i].startswith("    ") or lines[i].startswith("\t")
                    ):
                        collected.append(lines[i].lstrip())
                        i += 1
                    data[key] = "\n".join(collected).rstrip()
                    continue
                else:
                    data[key] = val
            else:
                # ignore unknown lines
                pass
            i += 1
        blocks.append(data)
    return blocks


def resolve_image_path(content: str):
    if not content:
        return None
    p = ROOT / content
    if p.exists():
        return p
    src = Path(content)
    parent = ROOT / src.parent
    stem = src.stem
    for ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp"):
        cand = parent / (stem + ext)
        if cand.exists():
            return cand
    for ext in (".png", ".jpg", ".jpeg"):
        cand = ROOT / "assets" / (stem + ext)
        if cand.exists():
            return cand
    return None


def apply_background(slide, img_path: Path, prs: Presentation):
    """Add a full-bleed picture as the slide background (placed behind other shapes)."""
    if not img_path or not img_path.exists():
        return
    try:
        pic = slide.shapes.add_picture(
            str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
        try:
            spTree = slide.shapes._spTree
            spTree.remove(pic._element)
            spTree.insert(2, pic._element)
        except Exception:
            pass
    except Exception:
        pass


def set_title(slide, title_text: str, font_color: str = None):
    if not title_text:
        return

    title_shape = None
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None

    if title_shape is None or not getattr(title_shape, "has_text_frame", False):
        cand = None
        min_top = None
        for shp in slide.shapes:
            if not getattr(shp, "has_text_frame", False):
                continue
            try:
                name = (shp.name or "").lower()
            except Exception:
                name = ""
            if "title" in name:
                cand = shp
                break
            try:
                t = getattr(shp, "top", None)
                if t is not None:
                    if min_top is None or t < min_top:
                        min_top = t
                        cand = shp
            except Exception:
                pass
        title_shape = cand

    if title_shape is None:
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False):
                title_shape = shp
                break

    if title_shape is None:
        return

    parsed_color = parse_hex_color(font_color)

    try:
        tf = title_shape.text_frame
        tf.clear()
        try:
            p = tf.paragraphs[0]
            # remove existing runs
            while len(p.runs) > 0:
                p.text = ""
                break
            r = p.add_run()
            r.text = title_text
            try:
                r.font.size = Pt(32)
                r.font.bold = True
                r.font.name = "Tenorite"
                if parsed_color is not None:
                    r.font.color.rgb = parsed_color
                else:
                    r.font.color.rgb = RGBColor(50, 50, 50)
            except Exception:
                pass
            try:
                p.alignment = PP_ALIGN.LEFT
            except Exception:
                pass
        except Exception:
            try:
                title_shape.text = title_text
            except Exception:
                pass
    except Exception:
        try:
            title_shape.text = title_text
        except Exception:
            pass


def set_body_text(slide, body_text: str, font_color: str = None):
    if not body_text:
        return
    parsed_color = parse_hex_color(font_color)
    for shp in slide.shapes:
        if shp.is_placeholder and shp.placeholder_format.type is not None:
            try:
                if shp == slide.shapes.title:
                    continue
            except Exception:
                pass
        if shp.has_text_frame and shp != slide.shapes.title:
            tf = shp.text_frame
            tf.clear()
            for i, line in enumerate(body_text.splitlines()):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = line.strip()
                p.font.size = Pt(18)
                try:
                    p.font.name = "Segoe UI"
                except Exception:
                    pass
                try:
                    # always set run font name; apply color if provided
                    for run in p.runs:
                        try:
                            run.font.name = "Segoe UI"
                        except Exception:
                            pass
                    if parsed_color is not None:
                        for run in p.runs:
                            run.font.color.rgb = parsed_color
                except Exception:
                    pass
            return
    tx = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    tf = tx.text_frame
    tf.clear()
    for i, line in enumerate(body_text.splitlines()):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line.strip()
        p.font.size = Pt(18)
        try:
            p.font.name = "Segoe UI"
        except Exception:
            pass
        try:
            for run in p.runs:
                try:
                    run.font.name = "Segoe UI"
                except Exception:
                    pass
            if parsed_color is not None:
                for run in p.runs:
                    run.font.color.rgb = parsed_color
        except Exception:
            pass


def add_two_content(
    slide,
    left_content: str,
    right_content: str,
    prs: Presentation,
    font_color: str = None,
):
    """Place left_content and right_content into left/right placeholders if available.
    """
    title_shape = None
    title_text = None
    title_ph_type = None
    try:
        title_shape = slide.shapes.title
        try:
            title_text = title_shape.text
        except Exception:
            title_text = None
        try:
            title_ph_type = title_shape.placeholder_format.type
        except Exception:
            title_ph_type = None
    except Exception:
        title_shape = None

    placeholders = []
    for shp in slide.shapes:
        if not getattr(shp, "is_placeholder", False):
            continue
        if title_shape is not None and shp is title_shape:
            continue
        try:
            ph_type = shp.placeholder_format.type
            if title_ph_type is not None and ph_type == title_ph_type:
                continue
        except Exception:
            pass
        try:
            if (
                title_text
                and shp.has_text_frame
                and shp.text.strip() == title_text.strip()
            ):
                continue
        except Exception:
            pass
        placeholders.append(shp)

    placeholders.sort(key=lambda s: getattr(s, "left", 0))
    left_ph = placeholders[0] if len(placeholders) >= 1 else None
    right_ph = placeholders[1] if len(placeholders) >= 2 else None

    def _populate_image_into_placeholder(ph, img_path: Path):
        if ph is None or img_path is None or not img_path.exists():
            return False
        try:
            left = ph.left
            top = ph.top
            width = ph.width
            height = ph.height
            try:
                spTree = slide.shapes._spTree
                spTree.remove(ph._element)
            except Exception:
                pass
            target_w = width
            target_h = height
            if PIL_AVAILABLE:
                try:
                    with Image.open(str(img_path)) as im:
                        img_w, img_h = im.size
                        img_ratio = img_w / img_h if img_h != 0 else 1.0
                        ph_ratio = width / height if height != 0 else 1.0
                        if ph_ratio > img_ratio:
                            target_h = height
                            target_w = int(target_h * img_ratio)
                        else:
                            target_w = width
                            target_h = int(target_w / img_ratio)
                except Exception:
                    target_w = width
                    target_h = height
            left_offset = left + max(0, (width - target_w) // 2)
            top_offset = top + max(0, (height - target_h) // 2)
            slide.shapes.add_picture(
                str(img_path), left_offset, top_offset, width=target_w, height=target_h
            )
            return True
        except Exception:
            return False

    parsed_color = parse_hex_color(font_color)

    def _populate_text_into_placeholder(ph, text: str, font_pt=14):
        if ph is None:
            return False
        try:
            if ph.has_text_frame:
                tf = ph.text_frame
                tf.clear()
                for i, line in enumerate(str(text).splitlines()):
                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    p.text = line.strip()
                    p.font.size = Pt(font_pt)
                    try:
                        p.font.name = "Segoe UI"
                    except Exception:
                        pass
                    try:
                        for run in p.runs:
                            try:
                                run.font.name = "Segoe UI"
                            except Exception:
                                pass
                        if parsed_color is not None:
                            for run in p.runs:
                                run.font.color.rgb = parsed_color
                    except Exception:
                        pass
                return True
        except Exception:
            return False
        return False

    def _maybe_image_path(item):
        if not item:
            return None
        r = resolve_image_path(str(item))
        if r:
            return r
        p = Path(item)
        return p if p.exists() else None

    left_img = _maybe_image_path(left_content)
    right_img = _maybe_image_path(right_content)

    left_done = False
    right_done = False

    if left_img:
        left_done = _populate_image_into_placeholder(left_ph, left_img)
    else:
        left_done = _populate_text_into_placeholder(left_ph, left_content, font_pt=14)

    if right_img:
        right_done = _populate_image_into_placeholder(right_ph, right_img)
    else:
        right_done = _populate_text_into_placeholder(
            right_ph, right_content, font_pt=14
        )

    if not left_done:
        tx = slide.shapes.add_textbox(Inches(0.3), Inches(1.5), Inches(4.5), Inches(4))
        tf = tx.text_frame
        tf.clear()
        for i, line in enumerate(str(left_content).splitlines()):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = line.strip()
            p.font.size = Pt(14)
            try:
                p.font.name = "Segoe UI"
            except Exception:
                pass
            try:
                if parsed_color is not None:
                    for run in p.runs:
                        run.font.color.rgb = parsed_color
                        try:
                            run.font.name = "Segoe UI"
                        except Exception:
                            pass
            except Exception:
                pass

    if not right_done:
        tx = slide.shapes.add_textbox(Inches(4.8), Inches(1.5), Inches(4.5), Inches(4))
        tf = tx.text_frame
        tf.clear()
        for i, line in enumerate(str(right_content).splitlines()):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = line.strip()
            p.font.size = Pt(14)
            try:
                p.font.name = "Segoe UI"
            except Exception:
                pass
            try:
                if parsed_color is not None:
                    for run in p.runs:
                        run.font.color.rgb = parsed_color
                        try:
                            run.font.name = "Segoe UI"
                        except Exception:
                            pass
            except Exception:
                pass


def find_layout(prs: Presentation, name_hint: str):
    if not name_hint:
        return prs.slide_layouts[6]
    name_hint = name_hint.strip().lower()
    for layout in prs.slide_layouts:
        try:
            if layout.name and name_hint in layout.name.lower():
                return layout
        except Exception:
            pass
    return prs.slide_layouts[6]


def main():
    if not CONTEXT.exists():
        print(f"context.md not found at {CONTEXT}")
        return
    text = CONTEXT.read_text(encoding="utf-8")
    blocks = parse_blocks(text)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    for b in blocks:
        layout_name = b.get("slide_layout", "")
        layout = find_layout(prs, layout_name)
        slide = prs.slides.add_slide(layout)

        bg = b.get("slide_bg")
        if bg:
            img = resolve_image_path(bg) or (ROOT / bg)
            apply_background(slide, img, prs)

        title = b.get("title")
        set_title(slide, title, b.get("font_color"))

        if "content1" in b or "content2" in b:
            left = b.get("content1", "")
            right = b.get("content2", "")
            add_two_content(slide, left, right, prs, b.get("font_color"))
        elif "content" in b:
            set_body_text(slide, b.get("content", ""), b.get("font_color"))

    prs.save(str(OUT_FILE))
    print(f"Saved presentation to {OUT_FILE}")


if __name__ == "__main__":
    main()
